from flask import Flask, request, jsonify
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from werkzeug.utils import secure_filename
import os

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

UPLOAD_FOLDER = "output_images"  # Change this to your preferred folder
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def extract_pptx_data(file_path):
    prs = Presentation(file_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        slide_info = {
            'slide_number': idx + 1,
            'shapes': []
        }

        for shape in slide.shapes:
            shape_data = {
                'x': shape.left,
                'y': shape.top,
                'width': shape.width,
                'height': shape.height,
                'type': 'UNKNOWN_SHAPE_TYPE'
            }

            # Set general shape type
            try:
                shape_data['type'] = MSO_SHAPE_TYPE(shape.shape_type).name
            except Exception:
                shape_data['type'] = f"UNKNOWN ({shape.shape_type})"

            # Handle AutoShapes (like Rectangle, Oval, Chevron, etc.)
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if hasattr(shape, "auto_shape_type"):
                    shape_data['auto_shape_type'] = shape.auto_shape_type
                    try:
                        shape_data['auto_shape_type_name'] = MSO_AUTO_SHAPE_TYPE(shape.auto_shape_type).name
                    except ValueError:
                        shape_data['auto_shape_type_name'] = f"UNKNOWN_AUTO_SHAPE ({shape.auto_shape_type})"

            # Extract text
            if shape.has_text_frame:
                shape_data['text'] = shape.text.strip()

            # Handle images
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, 'image'):
                image = shape.image
                ext = image.ext
                image_filename = f"slide_{idx + 1}_image.{ext}"
                image_path = os.path.join(UPLOAD_FOLDER, image_filename)

                with open(image_path, 'wb') as f:
                    f.write(image.blob)

                shape_data['image_filename'] = image_filename

            # Handle tables
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE and shape.has_table:
                shape_data['table'] = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]

            # You can add other types like LINE, GROUP, CHART, etc., if needed.

            slide_info['shapes'].append(shape_data)

        slides_data.append(slide_info)

    return slides_data



@app.route('/upload', methods=['POST'])
def upload_pptx():
    if 'pptx' not in request.files:
        return jsonify({'error': 'No PPTX file uploaded'}), 400

    file = request.files['pptx']
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)

    try:
        data = extract_pptx_data(filepath)
        return jsonify({'slides': data})
    except Exception as e:
        print(f"❌ Error: {e}")
        return jsonify({'error': str(e)}), 500
    finally:
        if os.path.exists(filepath):
            os.remove(filepath)


if __name__ == '__main__':
    app.run(port=5000)
